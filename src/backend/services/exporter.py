"""
슬라이드 PNG + 주석 래스터라이징 → PDF 내보내기

전략:
- PDF 원본: original.pdf 직접 복사 → 텍스트/검색/하이퍼링크 완전 유지 (OCR 지원)
- PPTX 원본: reportlab + python-pptx invisible text overlay → searchable PDF
  * 슬라이드 PNG를 배경 이미지로, python-pptx로 추출한 텍스트를 render_mode=3(invisible)으로 overlay
  * 위치는 EMU → px 스케일 변환으로 근사 매핑
- 유인물 레이아웃: Pillow 이미지 기반 (1up/2up/4up) — A4 기준, 노트 텍스트 포함
"""
from __future__ import annotations

import json
import logging
import platform
import re
import shutil
import textwrap
import urllib.request
from pathlib import Path
from typing import Literal

from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# A4 @ 150 DPI
A4_W = 1240  # px (210mm)
A4_H = 1754  # px (297mm)
MARGIN = 60   # px
NOTE_LINE_H = 28  # px per line
NOTE_FONT_SIZE = 18
SLIDE_NOTE_GAP = 16

_KOREAN_FONT_NAME: str | None = None
_KOREAN_FONT_CHECKED = False

# (우선순위, 등록명, 파일명 정규식(소문자), TTC subfont index)
_KOREAN_FONT_PRIORITY: list[tuple[int, str, str, int]] = [
    (0, "MalgunGothic",   r"^malgun\.ttf$",              0),
    (1, "MalgunGothicBd", r"^malgunbd\.ttf$",            0),
    (2, "NanumGothic",    r"^nanumgothic\.ttf$",         0),
    (3, "NanumMyeongjo",  r"^nanummyeongjo\.ttf$",       0),
    (4, "Batang",         r"^batang\.ttc$",              0),
    (5, "Gulim",          r"^gulim\.ttc$",               0),
    (6, "Dotum",          r"^dotum\.ttc$",               0),
    (7, "NotoSansCJK",    r"notosanscjk.*\.(ttf|ttc|otf)$", 0),
]


def _find_korean_font_paths() -> list[tuple[str, str, int]]:
    """시스템 폰트 디렉토리를 탐색해 한글 폰트 (등록명, 절대경로, ttc_index) 목록을 우선순위 순으로 반환."""
    sys_name = platform.system()
    if sys_name == "Windows":
        font_dirs = [
            Path(r"C:\Windows\Fonts"),
            Path.home() / "AppData" / "Local" / "Microsoft" / "Windows" / "Fonts",
        ]
    elif sys_name == "Darwin":
        font_dirs = [
            Path("/Library/Fonts"),
            Path("/System/Library/Fonts"),
            Path.home() / "Library" / "Fonts",
        ]
    else:
        font_dirs = [
            Path("/usr/share/fonts"),
            Path("/usr/local/share/fonts"),
            Path.home() / ".fonts",
            Path.home() / ".local" / "share" / "fonts",
        ]

    # 폰트 파일 수집 (Windows는 flat 디렉토리, 나머지는 재귀)
    all_fonts: list[Path] = []
    exts = ("*.ttf", "*.ttc", "*.otf")
    for d in font_dirs:
        if not d.exists():
            continue
        for ext in exts:
            if sys_name == "Windows":
                all_fonts.extend(d.glob(ext))
            else:
                all_fonts.extend(d.rglob(ext))

    # 중복 경로 제거
    seen_paths: set[str] = set()
    unique_fonts: list[Path] = []
    for f in all_fonts:
        key = str(f.resolve()).lower()
        if key not in seen_paths:
            seen_paths.add(key)
            unique_fonts.append(f)

    logger.debug("시스템 폰트 파일 탐색: %d개", len(unique_fonts))

    # 우선순위 패턴과 매칭
    matched: list[tuple[int, str, str, int]] = []
    seen_names: set[str] = set()
    for font_path in unique_fonts:
        fname_lower = font_path.name.lower()
        for priority, name, pattern, ttc_idx in _KOREAN_FONT_PRIORITY:
            if name not in seen_names and re.search(pattern, fname_lower):
                matched.append((priority, name, str(font_path), ttc_idx))
                seen_names.add(name)
                break

    matched.sort(key=lambda x: x[0])
    result = [(name, path, idx) for _, name, path, idx in matched]
    logger.info("한글 폰트 후보 발견 (%d개): %s", len(result), [(n, p) for n, p, _ in result])
    return result


def _get_korean_font_name() -> str | None:
    """시스템 폰트 탐색 → reportlab 등록 → 폰트 이름 반환. 없으면 None."""
    global _KOREAN_FONT_NAME, _KOREAN_FONT_CHECKED
    if _KOREAN_FONT_CHECKED:
        return _KOREAN_FONT_NAME
    _KOREAN_FONT_CHECKED = True
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        candidates = _find_korean_font_paths()
        if not candidates:
            logger.warning("시스템에서 한글 폰트를 찾지 못함 — invisible text 레이어 미적용")
            return None

        for name, path, ttc_idx in candidates:
            try:
                ext = Path(path).suffix.lower()
                if ext == ".ttc":
                    pdfmetrics.registerFont(TTFont(name, path, subfontIndex=ttc_idx))
                else:
                    pdfmetrics.registerFont(TTFont(name, path))
                _KOREAN_FONT_NAME = name
                logger.info("한글 폰트 등록 성공: %s (%s)", name, path)
                return name
            except Exception as e:
                logger.debug("폰트 등록 실패 %s: %s", path, e)
                continue

        logger.warning("한글 폰트 등록 실패 — invisible text 레이어 미적용")
    except ImportError:
        logger.warning("reportlab 없음 — PDF 내보내기 불가")
    return None


# ── PPTX 폰트 자동 탐색 + 인터넷 다운로드 ────────────────────────────────────────────────

_FONT_CACHE_DIR: Path = Path.home() / ".slidenote" / "fonts"

# PPTX에서 읽히는 다양한 폰트명 → 정규화된 등록명
_FONT_NAME_NORMALIZE: dict[str, str] = {
    "맑은 고딕": "MalgunGothic",       "맑은고딕": "MalgunGothic",
    "Malgun Gothic": "MalgunGothic",
    "나눔고딕": "NanumGothic",          "Nanum Gothic": "NanumGothic",
    "나눔명조": "NanumMyeongjo",        "Nanum Myeongjo": "NanumMyeongjo",
    "나눔바른고딕": "NanumBarunGothic", "Nanum Barun Gothic": "NanumBarunGothic",
    "나눔스퀘어": "NanumSquare",        "Nanum Square": "NanumSquare",
    "나눔손글씨붓": "NanumBrushScript", "Nanum Brush Script": "NanumBrushScript",
    "나눔손글씨펜": "NanumPenScript",   "Nanum Pen Script": "NanumPenScript",
    "돋움": "Dotum",    "굴림": "Gulim",    "바탕": "Batang",    "궁서": "Gungsuh",
    "Do Hyeon": "DoHyeon",   "도현": "DoHyeon",
    "Jua": "Jua",            "주아": "Jua",
    "Black Han Sans": "BlackHanSans",
    "Noto Sans KR": "NotoSansKR",      "노토산스KR": "NotoSansKR",
    "Noto Serif KR": "NotoSerifKR",
    "IBM Plex Sans KR": "IBMPlexSansKR",
    "Gowun Batang": "GowunBatang",     "고운바탕": "GowunBatang",
    "Gowun Dodum": "GowunDodum",       "고운돋움": "GowunDodum",
    "Song Myung": "SongMyung",         "송명": "SongMyung",
    "Yeon Sung": "YeonSung",           "연성": "YeonSung",
    "Gugi": "Gugi",   "Gaegu": "Gaegu",   "Poor Story": "PoorStory",
    "Hi Melody": "HiMelody",
    "East Sea Dokdo": "EastSeaDokdo",
    "Kirang Haerang": "KirangHaerang",
}

# 정규화 등록명 → Google Fonts GitHub 직접 TTF 다운로드 URL
_GOOGLE_FONTS_TTF_URLS: dict[str, str] = {
    "NanumGothic":      "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumgothic/NanumGothic-Regular.ttf",
    "NanumMyeongjo":    "https://raw.githubusercontent.com/google/fonts/main/ofl/nanummyeongjo/NanumMyeongjo-Regular.ttf",
    "NanumBarunGothic": "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumbarungothic/NanumBarunGothic.ttf",
    "NanumBrushScript": "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumbrushscript/NanumBrushScript-Regular.ttf",
    "NanumPenScript":   "https://raw.githubusercontent.com/google/fonts/main/ofl/nanumpenscript/NanumPenScript-Regular.ttf",
    "DoHyeon":          "https://raw.githubusercontent.com/google/fonts/main/ofl/dohyeon/DoHyeon-Regular.ttf",
    "Jua":              "https://raw.githubusercontent.com/google/fonts/main/ofl/jua/Jua-Regular.ttf",
    "BlackHanSans":     "https://raw.githubusercontent.com/google/fonts/main/ofl/blackhansans/BlackHanSans-Regular.ttf",
    "GowunBatang":      "https://raw.githubusercontent.com/google/fonts/main/ofl/gowunbatang/GowunBatang-Regular.ttf",
    "GowunDodum":       "https://raw.githubusercontent.com/google/fonts/main/ofl/gowundodum/GowunDodum-Regular.ttf",
    "IBMPlexSansKR":    "https://raw.githubusercontent.com/google/fonts/main/ofl/ibmplexsanskr/IBMPlexSansKR-Regular.ttf",
    "SongMyung":        "https://raw.githubusercontent.com/google/fonts/main/ofl/songmyung/SongMyung-Regular.ttf",
    "YeonSung":         "https://raw.githubusercontent.com/google/fonts/main/ofl/yeonsung/YeonSung-Regular.ttf",
    "Gugi":             "https://raw.githubusercontent.com/google/fonts/main/ofl/gugi/Gugi-Regular.ttf",
    "Gaegu":            "https://raw.githubusercontent.com/google/fonts/main/ofl/gaegu/Gaegu-Regular.ttf",
    "PoorStory":        "https://raw.githubusercontent.com/google/fonts/main/ofl/poorstory/PoorStory-Regular.ttf",
    "HiMelody":         "https://raw.githubusercontent.com/google/fonts/main/ofl/himelody/HiMelody-Regular.ttf",
    "KirangHaerang":    "https://raw.githubusercontent.com/google/fonts/main/ofl/kiranghaerang/KirangHaerang-Regular.ttf",
}


def _get_pptx_font_names(pptx_path: Path) -> set[str]:
    """PPTX 파일에서 사용된 폰트 이름 목록 반환 (단락·런 레벨 모두 포함)."""
    from pptx import Presentation
    fonts: set[str] = set()
    try:
        prs = Presentation(str(pptx_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                for para in shape.text_frame.paragraphs:
                    if para.font.name:
                        fonts.add(para.font.name)
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
    except Exception as e:
        logger.warning("PPTX 폰트 이름 추출 실패: %s", e)
    return fonts


def _try_download_font(reg_name: str) -> Path | None:
    """Google Fonts GitHub에서 TTF 파일 다운로드 → 캐시 경로 반환. 실패 시 None."""
    url = _GOOGLE_FONTS_TTF_URLS.get(reg_name)
    if not url:
        return None

    _FONT_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    ext = Path(url.split("?")[0]).suffix or ".ttf"
    cached = _FONT_CACHE_DIR / f"{reg_name}{ext}"
    if cached.exists():
        logger.debug("폰트 캐시 사용: %s", cached)
        return cached

    logger.info("폰트 다운로드 시작: %s → %s", reg_name, url)
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0 SlideNote/1.0"})
        with urllib.request.urlopen(req, timeout=20) as resp:
            data = resp.read()
        cached.write_bytes(data)
        logger.info("폰트 다운로드 완료: %s (%d bytes)", cached.name, len(data))
        return cached
    except Exception as e:
        logger.warning("폰트 다운로드 실패 (%s): %s", reg_name, e)
        cached.unlink(missing_ok=True)
        return None


def _build_font_registry(pptx_path: Path) -> dict[str, str]:
    """PPTX 폰트 목록 → {pptx_font_name: reportlab_reg_name} 매핑.

    순서: ① 이미 등록된 폰트 재사용 → ② 시스템 탐색 → ③ 인터넷 다운로드
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    pptx_fonts = _get_pptx_font_names(pptx_path)
    logger.info("PPTX 사용 폰트 목록: %s", pptx_fonts)

    sys_fonts = {name: (path, idx) for name, path, idx in _find_korean_font_paths()}
    registry: dict[str, str] = {}

    for font_name in pptx_fonts:
        canonical = _FONT_NAME_NORMALIZE.get(font_name, font_name.replace(" ", ""))

        # ① 이미 reportlab에 등록된 폰트 재사용
        try:
            pdfmetrics.getFont(canonical)
            registry[font_name] = canonical
            logger.debug("기등록 폰트 재사용: %s → %s", font_name, canonical)
            continue
        except Exception:
            pass

        # ② 시스템 폰트에서 찾기
        if canonical in sys_fonts:
            path, idx = sys_fonts[canonical]
            try:
                ext = Path(path).suffix.lower()
                if ext == ".ttc":
                    pdfmetrics.registerFont(TTFont(canonical, path, subfontIndex=idx))
                else:
                    pdfmetrics.registerFont(TTFont(canonical, path))
                registry[font_name] = canonical
                logger.info("시스템 폰트 등록: %s → %s (%s)", font_name, canonical, path)
                continue
            except Exception as e:
                logger.debug("시스템 폰트 등록 실패 %s: %s", font_name, e)

        # ③ 인터넷 다운로드 시도
        font_path = _try_download_font(canonical)
        if font_path:
            try:
                pdfmetrics.registerFont(TTFont(canonical, str(font_path)))
                registry[font_name] = canonical
                logger.info("다운로드 폰트 등록: %s → %s", font_name, canonical)
            except Exception as e:
                logger.warning("다운로드 폰트 등록 실패 %s: %s", font_name, e)

    logger.info("폰트 레지스트리 완성 (%d/%d): %s",
                len(registry), len(pptx_fonts), registry)
    return registry


def _get_shape_font_name(
    shape,
    font_registry: dict[str, str],
    fallback: str | None,
) -> str | None:
    """shape의 단락·런에서 사용된 폰트를 레지스트리에서 찾아 반환. 없으면 fallback."""
    if not hasattr(shape, "text_frame"):
        return fallback
    for para in shape.text_frame.paragraphs:
        if para.font.name and para.font.name in font_registry:
            return font_registry[para.font.name]
        for run in para.runs:
            if run.font.name and run.font.name in font_registry:
                return font_registry[run.font.name]
    return fallback


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    """시스템 한글 폰트 시도, 없으면 기본 폰트."""
    candidates = [
        "malgun.ttf",           # Windows 맑은 고딕
        "NanumGothic.ttf",      # Linux NanumGothic
        "NotoSansCJK-Regular.ttc",
    ]
    for name in candidates:
        try:
            return ImageFont.truetype(name, size)
        except OSError:
            continue
    return ImageFont.load_default()


def _draw_note_lines(draw: ImageDraw.ImageDraw, x: int, y: int, w: int,
                     text: str, font, line_h: int, max_lines: int = 6):
    """노트 텍스트를 지정 영역에 줄 바꿈해 그리기. 빈 줄 구분선도 표시."""
    lines: list[str] = []
    if text:
        for raw in text.splitlines():
            wrapped = textwrap.wrap(raw or " ", width=max(1, w // (NOTE_FONT_SIZE // 2 + 2)))
            lines.extend(wrapped or [""])
    # 빈 줄 채우기 (노트 필기 공간)
    while len(lines) < max_lines:
        lines.append("")

    for i, line in enumerate(lines[:max_lines]):
        ty = y + i * line_h
        # 연한 줄 그리기
        draw.line([(x, ty + line_h - 4), (x + w, ty + line_h - 4)], fill="#cccccc", width=1)
        if line:
            draw.text((x, ty), line, fill="#222222", font=font)


def _resize_slide(img: Image.Image, max_w: int, max_h: int) -> Image.Image:
    img = img.convert("RGB")
    img.thumbnail((max_w, max_h), Image.LANCZOS)
    return img


def _export_pdf_original(file_dir: Path, out_path: Path) -> Path:
    """PDF 원본 복사 → 텍스트/검색/하이퍼링크 레이어 완전 유지."""
    original = file_dir / "original.pdf"
    if not original.exists():
        raise FileNotFoundError(f"원본 PDF 없음: {original}")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(str(original), str(out_path))
    logger.info("PDF 원본 복사 완료: %s", out_path)
    return out_path


def _export_pptx_searchable(file_dir: Path, out_path: Path) -> Path:
    """PPTX: 슬라이드 PNG 배경 + python-pptx invisible text overlay → searchable PDF."""
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas as rl_canvas

    slides_dir = file_dir / "slides"
    png_files = sorted(slides_dir.glob("page_*.png"))
    if not png_files:
        raise FileNotFoundError(f"슬라이드 이미지 없음: {slides_dir}")

    # python-pptx 텍스트 추출
    prs = None
    original_pptx = file_dir / "original.pptx"
    if original_pptx.exists():
        try:
            from pptx import Presentation
            prs = Presentation(str(original_pptx))
        except Exception as e:
            logger.warning("python-pptx 로드 실패: %s", e)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    # 폰트 레지스트리: PPTX 원본 폰트 → reportlab 등록명 (시스템 우선, 없으면 다운로드)
    fallback_font = _get_korean_font_name()
    font_registry: dict[str, str] = {}
    if original_pptx.exists() and prs is not None:
        try:
            font_registry = _build_font_registry(original_pptx)
        except Exception as e:
            logger.warning("폰트 레지스트리 구축 실패: %s", e)

    c = rl_canvas.Canvas(str(out_path))
    for slide_idx, png_path in enumerate(png_files):
        img = Image.open(png_path)
        img_w, img_h = img.size
        c.setPageSize((img_w, img_h))

        # 배경: 슬라이드 PNG
        c.drawImage(ImageReader(str(png_path)), 0, 0, img_w, img_h)

        # Invisible text overlay (shape별 원본 폰트 적용)
        if prs is not None and slide_idx < len(prs.slides):
            slide = prs.slides[slide_idx]
            pptx_w = prs.slide_width   # EMU
            pptx_h = prs.slide_height  # EMU
            if pptx_w > 0 and pptx_h > 0:
                scale_x = img_w / pptx_w
                scale_y = img_h / pptx_h
                for shape in slide.shapes:
                    if not hasattr(shape, "text_frame"):
                        continue
                    tf = shape.text_frame
                    if not tf.text.strip():
                        continue
                    # shape별 폰트: PPTX 원본 폰트 → fallback (시스템 한글폰트)
                    shape_font = _get_shape_font_name(shape, font_registry, fallback_font)
                    if not shape_font:
                        continue
                    # shape 위치 → PDF 좌표 (PDF y=0이 하단 → 상단 기준으로 변환)
                    sx = float(shape.left or 0) * scale_x
                    sy_top = float(shape.top or 0) * scale_y
                    sh = float(shape.height or 1) * scale_y
                    pdf_y = img_h - sy_top
                    para_count = max(len(tf.paragraphs), 1)
                    font_size = max(6, int(sh / para_count * 0.65))
                    c.saveState()
                    c.setFont(shape_font, font_size)
                    t = c.beginText(sx, pdf_y)
                    t.setTextRenderMode(3)  # invisible: 검색·복사 가능, 화면에 미표시
                    for para in tf.paragraphs:
                        if not para.text:
                            continue
                        try:
                            t.textLine(para.text)
                        except (UnicodeEncodeError, Exception) as _te:
                            logger.debug("textLine 스킵 (슬라이드 %d): %s", slide_idx + 1, _te)
                    c.drawText(t)
                    c.restoreState()

        c.showPage()

    c.save()
    logger.info("PPTX searchable PDF 생성 완료: %s (%d슬라이드)", out_path, len(png_files))
    return out_path


def _export_images_pdf(slides_dir: Path, out_path: Path) -> Path:
    """폴백: PNG 이미지들을 Pillow로 합쳐 PDF (텍스트 레이어 없음)."""
    png_files = sorted(slides_dir.glob("page_*.png"))
    if not png_files:
        raise FileNotFoundError(f"슬라이드 이미지 없음: {slides_dir}")
    images = [Image.open(p).convert("RGB") for p in png_files]
    first, rest = images[0], images[1:]
    out_path.parent.mkdir(parents=True, exist_ok=True)
    first.save(str(out_path), "PDF", save_all=True, append_images=rest, resolution=150)
    return out_path


def export_to_pdf(file_dir: Path, out_path: Path) -> Path:
    """
    검색 가능한(searchable) PDF 내보내기.
    - PDF 원본: original.pdf 직접 복사 → 텍스트/검색 레이어 완전 유지
    - PPTX 원본: 슬라이드 PNG + python-pptx invisible text overlay (reportlab)
    반환: out_path
    """
    meta_path = file_dir / "metadata.json"
    ext = ""
    if meta_path.exists():
        try:
            ext = json.loads(meta_path.read_text(encoding="utf-8")).get("ext", "")
        except Exception:
            pass

    if ext == ".pdf":
        return _export_pdf_original(file_dir, out_path)
    elif ext == ".pptx":
        return _export_pptx_searchable(file_dir, out_path)
    else:
        return _export_images_pdf(file_dir / "slides", out_path)


def export_handout(
    file_dir: Path,
    out_path: Path,
    layout: Literal["1up", "2up", "4up"] = "2up",
) -> Path:
    """
    유인물 레이아웃 PDF 내보내기.

    layout="1up": A4 1페이지에 슬라이드 1장 + 노트 영역 6줄
    layout="2up": A4 1페이지에 슬라이드 2장 + 각 노트 영역 4줄
    layout="4up": A4 1페이지에 슬라이드 4장 (2×2 그리드, 요약 인쇄용)
    """
    slides_dir = file_dir / "slides"
    notes_dir = file_dir / "notes"

    png_files = sorted(slides_dir.glob("page_*.png"))
    if not png_files:
        raise FileNotFoundError(f"슬라이드 이미지 없음: {slides_dir}")

    font = _load_font(NOTE_FONT_SIZE)
    pages_per_sheet = {"1up": 1, "2up": 2, "4up": 4}[layout]
    slide_count = len(png_files)

    # 페이지별 노트 텍스트 미리 로드
    note_texts: list[str] = []
    for idx in range(slide_count):
        page_num = idx + 1
        note_path = notes_dir / f"slide_{page_num:02d}.json"
        if note_path.exists():
            try:
                data = json.loads(note_path.read_text(encoding="utf-8"))
                note_texts.append(data.get("text", ""))
            except Exception:
                note_texts.append("")
        else:
            note_texts.append("")

    a4_pages: list[Image.Image] = []

    if layout == "1up":
        slide_w = A4_W - 2 * MARGIN
        slide_max_h = int(A4_H * 0.55)
        note_lines = 8
        for idx, png_path in enumerate(png_files):
            canvas = Image.new("RGB", (A4_W, A4_H), "white")
            draw = ImageDraw.Draw(canvas)
            slide_img = _resize_slide(Image.open(png_path), slide_w, slide_max_h)
            sx = MARGIN + (slide_w - slide_img.width) // 2
            canvas.paste(slide_img, (sx, MARGIN))
            note_y = MARGIN + slide_img.height + SLIDE_NOTE_GAP
            note_h = A4_H - note_y - MARGIN
            actual_lines = min(note_lines, note_h // NOTE_LINE_H)
            _draw_note_lines(draw, MARGIN, note_y, slide_w,
                             note_texts[idx], font, NOTE_LINE_H, actual_lines)
            a4_pages.append(canvas)

    elif layout == "2up":
        cell_w = A4_W - 2 * MARGIN
        cell_slide_h = int((A4_H - 2 * MARGIN - SLIDE_NOTE_GAP) * 0.40)
        note_lines = 4
        # 슬라이드 2장씩 한 페이지
        for sheet_i in range(0, slide_count, 2):
            canvas = Image.new("RGB", (A4_W, A4_H), "white")
            draw = ImageDraw.Draw(canvas)
            for slot, idx in enumerate([sheet_i, sheet_i + 1]):
                if idx >= slide_count:
                    break
                slot_top = MARGIN + slot * (A4_H - 2 * MARGIN) // 2
                slide_img = _resize_slide(Image.open(png_files[idx]), cell_w, cell_slide_h)
                sx = MARGIN + (cell_w - slide_img.width) // 2
                canvas.paste(slide_img, (sx, slot_top))
                note_y = slot_top + slide_img.height + SLIDE_NOTE_GAP // 2
                slot_bottom = MARGIN + (slot + 1) * (A4_H - 2 * MARGIN) // 2
                avail_h = slot_bottom - note_y - MARGIN // 2
                actual_lines = min(note_lines, max(1, avail_h // NOTE_LINE_H))
                _draw_note_lines(draw, MARGIN, note_y, cell_w,
                                 note_texts[idx], font, NOTE_LINE_H, actual_lines)
            # 슬롯 구분선
            mid_y = A4_H // 2
            draw.line([(MARGIN, mid_y), (A4_W - MARGIN, mid_y)], fill="#dddddd", width=2)
            a4_pages.append(canvas)

    else:  # 4up
        cols, rows = 2, 2
        cell_w = (A4_W - 2 * MARGIN - MARGIN) // cols
        cell_h = (A4_H - 2 * MARGIN - MARGIN) // rows
        for sheet_i in range(0, slide_count, 4):
            canvas = Image.new("RGB", (A4_W, A4_H), "white")
            draw = ImageDraw.Draw(canvas)
            for slot in range(4):
                idx = sheet_i + slot
                if idx >= slide_count:
                    break
                col = slot % cols
                row = slot // cols
                cx = MARGIN + col * (cell_w + MARGIN)
                cy = MARGIN + row * (cell_h + MARGIN)
                slide_img = _resize_slide(Image.open(png_files[idx]), cell_w, cell_h - 20)
                sx = cx + (cell_w - slide_img.width) // 2
                canvas.paste(slide_img, (sx, cy))
                # 슬라이드 번호
                draw.text((cx, cy + slide_img.height + 4),
                           f"슬라이드 {idx + 1}", fill="#888888", font=font)
            # 그리드 구분선
            mid_x = A4_W // 2
            mid_y = A4_H // 2
            draw.line([(mid_x, MARGIN), (mid_x, A4_H - MARGIN)], fill="#dddddd", width=1)
            draw.line([(MARGIN, mid_y), (A4_W - MARGIN, mid_y)], fill="#dddddd", width=1)
            a4_pages.append(canvas)

    if not a4_pages:
        raise RuntimeError("생성된 유인물 페이지 없음")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    first, rest = a4_pages[0], a4_pages[1:]
    first.save(str(out_path), "PDF", save_all=True, append_images=rest, resolution=150)
    return out_path


"""
Gemini Vision 슬라이드 요약 서비스

검증된 패턴 (gen_vision_improved.py 기반):
- PRIMARY: gemini-2.5-flash-image → inline_data 감지 시 FALLBACK: gemini-2.5-flash
- DELAY: 7.0초 (무료 티어 ~10 RPM 제한)
- 온도 0.1, max_tokens 8192
"""
from __future__ import annotations

import asyncio
import logging
import os
from pathlib import Path

logger = logging.getLogger(__name__)

GEMINI_PRIMARY = "gemini-2.5-flash-image"
GEMINI_FALLBACK = "gemini-2.5-flash"
GEMINI_DELAY = 7.0

SLIDE_PROMPT = """다음 슬라이드 이미지를 보고 **발표자 노트 형식**으로 한국어 요약을 작성해 주세요.

## 작성 형식

### 핵심 요약 (3줄 이내)
슬라이드의 핵심 메시지 3줄

### 주요 내용
- 슬라이드에 담긴 주요 내용을 불릿으로 정리
- 다이어그램, 흐름도, 표가 있으면 설명 포함

### 발표 포인트
발표자가 강조해야 할 한 가지 포인트

---
Markdown 형식으로만 출력하세요. 마크다운 코드블록(```) 래퍼 없이 바로 내용만 출력하세요."""

# ── 슬라이드 → Markdown 변환용 프롬프트 (AI 요약과 별개) ──────────────────────
# 목적: Obsidian / Notion / NotebookLM 임포트에 최적화된 원문 충실 변환
# AI 요약(SLIDE_PROMPT)과의 차이:
#   SLIDE_PROMPT  → 3줄 핵심 요약 + 발표 포인트 (내용 압축)
#   SLIDE_TO_MARKDOWN_PROMPT → 슬라이드에 있는 모든 텍스트/표/다이어그램을 빠짐없이 구조화
SLIDE_TO_MARKDOWN_PROMPT = """다음 슬라이드 이미지를 Markdown으로 변환하세요.
요약하거나 내용을 줄이지 마세요. 슬라이드에 있는 텍스트는 모두 포함해야 합니다.

## 변환 규칙

### 1. 헤딩 (시각적 위계 기반 동적 매핑)
- 슬라이드에서 **가장 크고 두드러진** 텍스트 하나 → `## 제목` (슬라이드당 `##` 헤딩은 하나만)
- 다음 단계 (섹션 번호·카테고리명) → `### 소제목`
- 더 세부적인 소항목 레이블 → `#### 세부항목`
- 헤딩 레벨 건너뛰기 금지 (`##` 바로 다음 `####`로 내려가지 말 것)
- 번호 레이블(01, 02 등)은 단독 헤딩 금지 → 뒤 텍스트와 합쳐서 같은 헤딩으로: `### 01 주요 성과`
- 동일 시각 크기의 텍스트 여러 개 → 모두 같은 레벨로 처리

### 2. 2컬럼 레이아웃
- 나란히 배치된 두 블록 → HTML table (컬럼 너비 50%)
- colspan/rowspan 있는 병합 셀 → HTML table with colspan/rowspan 속성 명시
- HTML `<td>` 안에서는 반드시 `<h3>`, `<h4>` 태그 사용 (`###`, `####` 등 Markdown 헤딩 절대 금지)
- HTML `<td>` 안에서 `style="..."` 인라인 스타일 사용 금지 → 강조는 `<strong>` 또는 일반 텍스트

### 3. 데이터 표
- 격자형 표 → Markdown 테이블 우선 (`| 헤더 | ... |`)
- 병합 셀 포함 시 → HTML `<table>` with `colspan`/`rowspan`

### 4. 다이어그램/흐름도/아키텍처
- 화살표·연결선·박스가 있는 도형 → Mermaid flowchart
  ```mermaid
  graph TD
    A[시작] --> B[처리] --> C[종료]
  ```
- 순서/절차 다이어그램 → Mermaid sequence
  ```mermaid
  sequenceDiagram
    Client->>Server: 요청
    Server-->>Client: 응답
  ```

### 5. 타임라인/간트/일정
- 날짜/기간/단계가 포함된 시각화 → Mermaid gantt
  ```mermaid
  gantt
    title 프로젝트 일정
    section 단계
    작업명 :a1, 2024-01-01, 30d
  ```
- 간단한 시간표 → Markdown 테이블

### 6. Step 카드 / 프로세스 번호 순서
- 번호 붙은 단계 블록 → HTML table (번호 bold) 또는 순서 목록

### 7. KEY TAKEAWAY / 인사이트 박스
- 강조 박스 → Markdown blockquote (`>`)

### 8. 색상 강조 텍스트
- 파란색/주요 키워드 → `**굵게**`

### 9. CLI/코드 예시
- 명령어, 코드 → ` ```bash ` 또는 ` ```sql ` 코드블록

### 10. 이미지/차트 (텍스트 전환 불가 시)
- 막대/파이/선 그래프에 **수치가 보이면** → 반드시 값을 담은 Markdown 테이블로 표현 (이미지 주석 처리 금지)
- 수치가 없는 사진·순수 이미지 → `<!-- 이미지: [설명] -->`

### 11. 페이지 번호, 회사 로고
- 생략

출력: Markdown만. 코드블록(```) 래퍼 없이 바로 내용 출력. 한국어 그대로 유지."""

# PPTX 구조 메타데이터 보강 프롬프트 (python-pptx 추출 결과를 주입)
_PPTX_META_PROMPT = """
---
## 슬라이드 구조 힌트 (python-pptx 추출)
{meta}
---
위 구조 힌트를 참고하여 위 규칙에 따라 슬라이드 이미지를 Markdown으로 변환하세요."""


def _get_client():
    """google-genai 클라이언트 초기화. API 키는 환경변수 GOOGLE_API_KEY."""
    from google import genai

    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        raise ValueError("GOOGLE_API_KEY 환경변수가 설정되지 않았습니다.")
    return genai.Client(api_key=api_key)


def _img_part(gtypes, img_path: Path):
    """이미지 파일을 google-genai Part로 변환."""
    img_bytes = img_path.read_bytes()
    suffix = img_path.suffix.lower().lstrip(".")
    mime = "image/png" if suffix == "png" else f"image/{suffix}"
    return gtypes.Part.from_bytes(data=img_bytes, mime_type=mime)


def _call_model(client, gtypes, model: str, img_path: Path, prompt: str = SLIDE_PROMPT) -> tuple[str, str | None]:
    """단일 모델 호출. 반환: (text, error_type|None)"""
    try:
        resp = client.models.generate_content(
            model=model,
            contents=[_img_part(gtypes, img_path), prompt],
            config=gtypes.GenerateContentConfig(
                temperature=0.1,
                max_output_tokens=8192,
            ),
        )
        text = resp.text or ""
        if not text.strip():
            return "", "inline_data"
        return text, None
    except Exception as e:
        err = str(e)
        if "inline_data" in err.lower() or "non-text" in err.lower():
            return "", "inline_data"
        return "", err


async def summarize_slide(img_path: Path) -> str:
    """
    슬라이드 PNG → 한국어 발표자 노트 요약 (Markdown).
    PRIMARY 실패 시 FALLBACK, GEMINI_DELAY 대기.
    """
    from google.genai import types as gtypes

    def _run() -> str:
        client = _get_client()
        text, err = _call_model(client, gtypes, GEMINI_PRIMARY, img_path)
        if err == "inline_data":
            logger.info("inline_data 감지 → %s 폴백", GEMINI_FALLBACK)
            import time
            time.sleep(GEMINI_DELAY)
            text, err = _call_model(client, gtypes, GEMINI_FALLBACK, img_path)
        if err:
            raise RuntimeError(f"Gemini 오류: {err}")
        return text

    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _run)


async def convert_slide_to_markdown(
    img_path: Path,
    pptx_path: Path | None = None,
    pdf_path: Path | None = None,
    slide_index: int = 0,
) -> str:
    """
    슬라이드 PNG → Obsidian/Notion/NotebookLM에 최적화된 구조화된 Markdown.

    AI 요약(summarize_slide)과 다름:
    - summarize_slide: 내용을 3줄로 압축하는 '발표자 노트 요약'
    - convert_slide_to_markdown: 슬라이드의 모든 내용을 빠짐없이 변환하는 '원문 충실 변환'

    pptx_path 또는 pdf_path가 주어지면 텍스트를 추출해 프롬프트에 보강한다
    (인식 정확도 향상). slide_index는 0-based.
    """
    from google.genai import types as gtypes

    # 소스 파일 메타데이터 추출 (동기, 빠름)
    extra_context = ""
    if pptx_path is not None and pptx_path.exists():
        try:
            extra_context = _extract_pptx_metadata(pptx_path, slide_index)
        except Exception as e:
            logger.warning("PPTX 메타데이터 추출 실패 (무시): %s", e)
    elif pdf_path is not None and pdf_path.exists():
        try:
            extra_context = _extract_pdf_metadata(pdf_path, slide_index)
        except Exception as e:
            logger.warning("PDF 메타데이터 추출 실패 (무시): %s", e)

    def _run() -> str:
        client = _get_client()

        # 메타데이터가 있으면 프롬프트 끝에 보강 컨텍스트 추가
        prompt = SLIDE_TO_MARKDOWN_PROMPT
        if extra_context:
            prompt = SLIDE_TO_MARKDOWN_PROMPT + _PPTX_META_PROMPT.format(meta=extra_context)

        text, err = _call_model(client, gtypes, GEMINI_PRIMARY, img_path, prompt)
        if err == "inline_data":
            logger.info("inline_data 감지 → %s 폴백 (Markdown 변환)", GEMINI_FALLBACK)
            import time
            time.sleep(GEMINI_DELAY)
            text, err = _call_model(client, gtypes, GEMINI_FALLBACK, img_path, prompt)
        if err:
            raise RuntimeError(f"Gemini Markdown 변환 오류: {err}")
        return text

    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _run)


def _extract_pptx_metadata(pptx_path: Path, slide_index: int) -> str:
    """python-pptx로 슬라이드 구조 텍스트 추출 (0-based index).

    텍스트박스·표 구조를 Gemini 프롬프트에 주입해 인식 정확도를 높인다.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides_list = list(prs.slides)
    if slide_index >= len(slides_list):
        return ""
    slide = slides_list[slide_index]
    lines: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(f"- {t}")
        elif shape.has_table:
            tbl = shape.table
            lines.append(f"[표: {len(tbl.rows)}행 x {len(tbl.columns)}열]")
            for row in tbl.rows:
                row_texts = [cell.text.strip() for cell in row.cells]
                lines.append("  | " + " | ".join(row_texts) + " |")
    return "\n".join(lines)


def _extract_pdf_metadata(pdf_path: Path, page_index: int) -> str:
    """PyMuPDF로 PDF 페이지 텍스트 추출 (0-based index).

    텍스트 블록을 Gemini 프롬프트에 주입해 인식 정확도를 높인다.
    스캔 PDF(텍스트 레이어 없음)는 빈 문자열 반환.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(str(pdf_path))
    if page_index >= len(doc):
        doc.close()
        return ""

    page = doc[page_index]
    # 텍스트 블록 추출: (x0, y0, x1, y1, text, block_no, block_type)
    blocks = page.get_text("blocks")
    doc.close()

    lines: list[str] = []
    for block in sorted(blocks, key=lambda b: (b[1], b[0])):  # y→x 순 정렬
        text = block[4].strip()
        if text:
            # 줄바꿈 정규화 (멀티라인 블록을 단일 불릿으로)
            single_line = " ".join(text.split())
            lines.append(f"- {single_line}")

    return "\n".join(lines)
